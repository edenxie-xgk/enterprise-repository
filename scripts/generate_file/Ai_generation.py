"""
多模态金融垂直数据生成系统 (Multimodal Financial Data Generator)
生成文本、图表、图像、表格、扫描件、语音转写等混合模态数据
"""

import os
import random
import json
import base64
import hashlib
import tempfile
from datetime import datetime, timedelta
from io import BytesIO
from typing import List, Dict, Optional, Tuple
import time

import pandas as pd
import numpy as np
from faker import Faker
from tqdm import tqdm
from PIL import Image, ImageDraw, ImageFont, ImageFilter, ImageEnhance
import matplotlib
matplotlib.use('Agg')  # 使用非交互式后端，避免GUI问题
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from matplotlib.patches import Rectangle, FancyBboxPatch
import seaborn as sns

# 文档处理
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from openpyxl import Workbook
from openpyxl.drawing.image import Image as XLImage
from openpyxl.chart import BarChart, Reference
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.enum.text import PP_ALIGN
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.graphics.shapes import Drawing
from reportlab.graphics.charts.linecharts import HorizontalLineChart
from reportlab.graphics.charts.barcharts import VerticalBarChart

# AI生成
from openai import OpenAI
import dotenv
from tenacity import retry, stop_after_attempt, wait_exponential

dotenv.load_dotenv()

# ============================================
# 配置与初始化
# ============================================

fake = Faker("zh_CN")
ROOT = "multimodal_financial_data"
os.makedirs(ROOT, exist_ok=True)

# 创建临时目录
TEMP_DIR = tempfile.mkdtemp(prefix="fin_data_")
print(f"临时目录: {TEMP_DIR}")

# 中文字体配置（多平台兼容）
FONT_PATHS = [
    "/System/Library/Fonts/PingFang.ttc",  # macOS
    "C:/Windows/Fonts/simhei.ttf",  # Windows
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",  # Linux
    "/fonts/NotoSansCJKsc-VF.ttf",  # Docker
]

CHINESE_FONT = None
for fp in FONT_PATHS:
    if os.path.exists(fp):
        CHINESE_FONT = fp
        break

# 如果没有找到中文字体，使用默认字体（可能不支持中文）
if not CHINESE_FONT:
    print("警告: 未找到中文字体，中文显示可能异常")

# 注册ReportLab字体
REPORTLAB_FONT = "Helvetica"
if CHINESE_FONT and CHINESE_FONT.endswith(('.ttf', '.ttc')):
    try:
        font_name = "CustomCN"
        pdfmetrics.registerFont(TTFont(font_name, CHINESE_FONT))
        REPORTLAB_FONT = font_name
    except Exception as e:
        print(f"字体注册失败: {e}")

# 颜色主题（专业金融配色）
COLORS = {
    'primary': '#1f4e79',  # 深蓝
    'secondary': '#2e75b6',  # 中蓝
    'accent': '#c55a11',  # 橙红（强调）
    'positive': '#70ad47',  # 绿色（上涨）
    'negative': '#c00000',  # 红色（下跌）
    'neutral': '#7f7f7f',  # 灰色
    'bg': '#f2f2f2',  # 背景灰
    'white': '#ffffff',
    'text': '#333333'
}

# 设置matplotlib中文字体
plt.rcParams['font.sans-serif'] = ['SimHei', 'DejaVu Sans', 'PingFang SC', 'Microsoft YaHei']
plt.rcParams['axes.unicode_minus'] = False


# ============================================
# 多模态内容生成器
# ============================================

class MultimodalContentEngine:
    """核心引擎：协调文本、图表、图像的生成"""

    def __init__(self):
        api_key = os.getenv("DEEPSEEK_API_KEY")
        if api_key:
            self.client = OpenAI(
                api_key=api_key,
                base_url="https://api.deepseek.com"
            )
        else:
            self.client = None
            print("警告: 未设置DEEPSEEK_API_KEY，将使用备用文本生成")

        self.companies = self._init_company_database()
        self.market_data = self._init_market_data()

    def _init_company_database(self) -> List[Dict]:
        """初始化真实感公司数据库"""
        sectors = {
            "半导体": ["中芯国际", "韦尔股份", "兆易创新", "北方华创", "紫光国微"],
            "新能源": ["宁德时代", "比亚迪", "隆基绿能", "通威股份", "阳光电源"],
            "医药生物": ["恒瑞医药", "迈瑞医疗", "药明康德", "爱尔眼科", "智飞生物"],
            "人工智能": ["科大讯飞", "海康威视", "商汤科技", "寒武纪", "云从科技"],
            "金融科技": ["东方财富", "同花顺", "恒生电子", "拉卡拉", "银之杰"]
        }

        companies = []
        for sector, names in sectors.items():
            for name in names:
                companies.append({
                    "name": name,
                    "sector": sector,
                    "ticker": self._generate_ticker(name),
                    "market_cap": random.randint(50, 800) * 1e8,  # 市值
                    "pe_ratio": round(random.uniform(15, 80), 2),
                    "pb_ratio": round(random.uniform(1.5, 8), 2),
                    "roe": round(random.uniform(5, 25), 2),
                    "revenue_growth": round(random.uniform(-10, 50), 2),
                    "profit_margin": round(random.uniform(5, 30), 2),
                    "founded": random.randint(1995, 2015),
                    "employees": random.randint(1000, 50000)
                })
        return companies

    def _generate_ticker(self, name: str) -> str:
        """生成股票代码"""
        hash_val = int(hashlib.md5(name.encode()).hexdigest(), 16)
        if "科技" in name or "智能" in name:
            prefix = ["688", "300"][hash_val % 2]  # 科创板/创业板
        else:
            prefix = ["600", "000", "002"][hash_val % 3]  # 主板/中小板
        return f"{prefix}{hash_val % 1000:03d}"

    def _init_market_data(self) -> Dict:
        """初始化市场基准数据"""
        dates = pd.date_range('2023-01-01', '2024-12-31', freq='B')  # 工作日
        np.random.seed(42)

        # 生成随机游走价格
        returns = np.random.normal(0.0005, 0.02, len(dates))
        prices = 3000 * np.exp(np.cumsum(returns))

        return {
            'dates': dates,
            'index_prices': prices,
            'volumes': np.random.randint(1e8, 5e8, len(dates))
        }

    @retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=2, max=10))
    def generate_analysis_text(self, context: Dict) -> str:
        """AI生成专业分析文本"""
        if not self.client:
            return self._fallback_analysis(context)

        prompt = f"""作为资深金融分析师，基于以下数据撰写{context.get('length', '200')}字的专业分析：

公司：{context.get('company', '行业整体')}
行业：{context.get('sector', '综合')}
关键指标：{json.dumps(context.get('metrics', {}), ensure_ascii=False)}

要求：
1. 结合具体数据，避免空泛描述
2. 包含趋势判断和投资建议
3. 使用专业术语（估值修复、业绩兑现、景气度等）
4. 直接输出正文，无标题和标记"""

        try:
            response = self.client.chat.completions.create(
                model="deepseek-chat",
                messages=[
                    {"role": "system", "content": "你是拥有10年经验的TMT行业首席分析师"},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                max_tokens=500
            )
            return response.choices[0].message.content.strip()
        except Exception as e:
            print(f"AI生成失败: {e}")
            return self._fallback_analysis(context)

    def _fallback_analysis(self, context: Dict) -> str:
        """备用文本"""
        templates = [
            "从估值角度看，当前PE处于历史中枢偏下位置，具备一定安全边际。",
            "业绩增速符合预期，毛利率环比改善明显，显示议价能力增强。",
            "行业景气度持续上行，龙头公司市占率提升逻辑清晰。",
            "短期受宏观情绪影响波动加大，但基本面支撑坚实。",
            "研发投入持续增加，新产品管线有望在未来两年贡献增量收入。",
            "海外市场拓展顺利，全球化布局初见成效。"
        ]
        return random.choice(templates)

    def generate_stock_chart(self, company: Dict, days: int = 180) -> Image.Image:
        """生成专业K线图/趋势图"""
        # 生成模拟股价数据
        end_date = datetime(2024, 12, 31)
        dates = pd.date_range(end=end_date, periods=days, freq='B')

        np.random.seed(hash(company['name']) % 2 ** 32)
        trend = random.choice(['up', 'down', 'volatile', 'stable'])

        # 根据趋势生成价格
        base_price = random.uniform(20, 200)
        if trend == 'up':
            returns = np.random.normal(0.002, 0.025, days)
        elif trend == 'down':
            returns = np.random.normal(-0.001, 0.025, days)
        elif trend == 'volatile':
            returns = np.random.normal(0, 0.04, days)
        else:
            returns = np.random.normal(0, 0.015, days)

        prices = base_price * np.exp(np.cumsum(returns))

        # 创建图表
        fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(10, 6),
                                       gridspec_kw={'height_ratios': [3, 1]},
                                       facecolor='white')

        # 价格线
        color = COLORS['positive'] if prices[-1] > prices[0] else COLORS['negative']
        ax1.plot(dates, prices, color=color, linewidth=2, label=company['name'])
        ax1.fill_between(dates, prices, alpha=0.3, color=color)

        # 均线
        ma20 = pd.Series(prices).rolling(20).mean()
        ma60 = pd.Series(prices).rolling(60).mean()
        ax1.plot(dates, ma20, '--', color=COLORS['secondary'], alpha=0.7, label='MA20')
        ax1.plot(dates, ma60, '--', color=COLORS['accent'], alpha=0.7, label='MA60')

        # 标注关键点位
        max_idx = np.argmax(prices)
        min_idx = np.argmin(prices)
        ax1.annotate(f'高点: ¥{prices[max_idx]:.2f}',
                     xy=(dates[max_idx], prices[max_idx]),
                     xytext=(10, 10), textcoords='offset points',
                     bbox=dict(boxstyle='round,pad=0.3', facecolor='yellow', alpha=0.7),
                     arrowprops=dict(arrowstyle='->', connectionstyle='arc3,rad=0'))

        ax1.set_title(f"{company['name']}({company['ticker']}) 股价走势", fontsize=14, fontweight='bold')
        ax1.legend(loc='upper left')
        ax1.grid(True, alpha=0.3)

        # 成交量
        volumes = np.random.randint(1e6, 5e7, days)
        colors_vol = [COLORS['positive'] if prices[i] > prices[i - 1] else COLORS['negative']
                      for i in range(1, days)]
        colors_vol.insert(0, COLORS['neutral'])
        ax2.bar(dates, volumes, color=colors_vol, alpha=0.7)
        ax2.set_ylabel('成交量')

        plt.tight_layout()

        # 转为PIL Image
        buf = BytesIO()
        plt.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
        buf.seek(0)
        img = Image.open(buf)
        plt.close(fig)  # 关闭图形释放内存

        return img

    def generate_financial_table_image(self, company: Dict) -> Image.Image:
        """生成财务报表截图样式（模拟Excel/系统截图）"""
        # 创建数据
        periods = ['2021A', '2022A', '2023A', '2024E', '2025E']
        metrics = {
            '营业收入(亿元)': [random.uniform(50, 200) * (1.15 ** i) for i in range(5)],
            '净利润(亿元)': [random.uniform(5, 30) * (1.2 ** i) for i in range(5)],
            '毛利率(%)': [random.uniform(25, 45) for _ in range(5)],
            '净利率(%)': [random.uniform(8, 20) for _ in range(5)],
            'ROE(%)': [random.uniform(10, 25) for _ in range(5)],
            '资产负债率(%)': [random.uniform(30, 60) for _ in range(5)]
        }

        # 创建图像
        fig, ax = plt.subplots(figsize=(12, 6))
        ax.axis('tight')
        ax.axis('off')

        # 表格数据
        table_data = [[metric] + [f'{val:.2f}' if isinstance(val, float) else str(val)
                                  for val in values]
                      for metric, values in metrics.items()]
        headers = ['指标'] + periods

        # 绘制表格
        table = ax.table(cellText=table_data, colLabels=headers,
                         cellLoc='center', loc='center',
                         colWidths=[0.2] + [0.16] * 5)

        table.auto_set_font_size(False)
        table.set_fontsize(10)
        table.scale(1, 2)

        # 样式设置
        for i in range(len(headers)):
            table[(0, i)].set_facecolor(COLORS['primary'])
            table[(0, i)].set_text_props(weight='bold', color='white')

        # 交替行颜色
        for i in range(1, len(table_data) + 1):
            for j in range(len(headers)):
                if i % 2 == 0:
                    table[(i, j)].set_facecolor('#f0f0f0')

        # 添加标题
        plt.title(f"{company['name']} 财务数据摘要", fontsize=16, fontweight='bold', pad=20)

        buf = BytesIO()
        plt.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
        buf.seek(0)
        img = Image.open(buf)
        plt.close(fig)

        return img

    def generate_scan_document(self, doc_type: str = "research_notes") -> Image.Image:
        """生成扫描件/手写笔记图像（带噪声、旋转、阴影）"""
        width, height = 1200, 1600
        img = Image.new('RGB', (width, height), color='white')
        draw = ImageDraw.Draw(img)

        # 添加横线（笔记本效果）
        line_spacing = 40
        for y in range(100, height, line_spacing):
            draw.line([(50, y), (width - 50, y)], fill='#e0e0e0', width=1)

        # 添加竖线（边栏）
        draw.line([(150, 50), (150, height - 50)], fill='#ffcccc', width=2)

        # 添加手写体文本（模拟）
        try:
            if CHINESE_FONT:
                font = ImageFont.truetype(CHINESE_FONT, 28)
                title_font = ImageFont.truetype(CHINESE_FONT, 36)
            else:
                font = ImageFont.load_default()
                title_font = font
        except:
            font = ImageFont.load_default()
            title_font = font

        # 标题
        title = f"{'调研' if doc_type == 'research_notes' else '会议纪要'} - {fake.date()}"
        draw.text((180, 80), title, fill='#333333', font=title_font)

        # 生成手写内容
        y_pos = 150
        content_lines = [
            f"调研对象：{random.choice(self.companies)['name']}",
            f"调研时间：{fake.date()} {fake.time()}",
            f"参与人员：{fake.name()}、{fake.name()}",
            "",
            "核心要点：",
            f"1. {self._fallback_analysis({})}",
            f"2. 订单情况：Q{random.randint(1, 4)}环比{'增长' if random.random() > 0.3 else '下滑'}{random.randint(5, 30)}%",
            f"3. 产能利用率：{random.randint(60, 95)}%",
            f"4. 竞争对手动态：{random.choice(['价格战趋缓', '新品发布延迟', '渠道扩张加速'])}",
            "",
            "投资建议：",
            random.choice(['维持买入', '上调至增持', '暂观望', '推荐关注'])
        ]

        for line in content_lines:
            draw.text((180, y_pos), line, fill='#2c5aa0' if random.random() > 0.1 else '#000080',
                      font=font)
            y_pos += line_spacing

        # 添加手写圈注（红色）
        draw.ellipse([(170, 400), (190, 420)], outline='red', width=2)
        draw.text((200, 405), "重点", fill='red', font=font)

        # 添加扫描件效果（噪声、模糊、旋转）
        img = img.filter(ImageFilter.GaussianBlur(radius=0.5))

        # 添加随机噪声
        np_img = np.array(img)
        noise = np.random.normal(0, 5, np_img.shape).astype(np.int16)
        np_img = np.clip(np_img.astype(np.int16) + noise, 0, 255).astype(np.uint8)
        img = Image.fromarray(np_img)

        # 轻微旋转（-2到2度）
        angle = random.uniform(-2, 2)
        img = img.rotate(angle, fillcolor='white', expand=False)

        # 添加阴影效果（模拟扫描光线不均）
        enhancer = ImageEnhance.Brightness(img)
        img = enhancer.enhance(random.uniform(0.9, 1.1))

        return img

    def generate_comparison_chart(self, companies: List[Dict]) -> Image.Image:
        """生成同业对比雷达图/柱状图"""
        metrics = ['PE', 'PB', 'ROE', '营收增速', '利润率', '负债率']

        fig = plt.figure(figsize=(14, 6))

        # 左图：雷达图
        ax1 = plt.subplot(121, projection='polar')
        angles = np.linspace(0, 2 * np.pi, len(metrics), endpoint=False).tolist()
        angles += angles[:1]

        colors_list = plt.cm.Set3(np.linspace(0, 1, len(companies)))

        for idx, comp in enumerate(companies[:3]):  # 最多对比3家
            # 归一化处理
            pe_norm = max(0, min(100, 100 - comp['pe_ratio']))
            pb_norm = max(0, min(100, 100 - comp['pb_ratio'] * 10))
            roe_norm = max(0, min(100, comp['roe'] * 4))
            growth_norm = max(0, min(100, comp['revenue_growth'] * 2 if comp['revenue_growth'] > 0 else 0))
            margin_norm = max(0, min(100, comp['profit_margin'] * 3))
            debt_norm = max(0, min(100, 100 - comp.get('debt_ratio', 50)))

            values = [pe_norm, pb_norm, roe_norm, growth_norm, margin_norm, debt_norm]
            values += values[:1]

            ax1.plot(angles, values, 'o-', linewidth=2, label=comp['name'], color=colors_list[idx])
            ax1.fill(angles, values, alpha=0.15, color=colors_list[idx])

        ax1.set_xticks(angles[:-1])
        ax1.set_xticklabels(metrics)
        ax1.set_ylim(0, 100)
        ax1.legend(loc='upper right', bbox_to_anchor=(1.3, 1.1))
        ax1.set_title('财务指标对比', fontsize=12, fontweight='bold', pad=20)

        # 右图：估值散点图
        ax2 = plt.subplot(122)
        x = [c['pe_ratio'] for c in companies]
        y = [c['pb_ratio'] for c in companies]
        sizes = [c['market_cap'] / 1e9 for c in companies]  # 市值决定大小
        colors_scatter = [COLORS['positive'] if c['revenue_growth'] > 20 else
                          COLORS['accent'] if c['revenue_growth'] > 0 else
                          COLORS['negative'] for c in companies]

        scatter = ax2.scatter(x, y, s=[s / 10 for s in sizes], c=colors_scatter, alpha=0.6)

        # 标注部分公司
        for comp in random.sample(companies, min(5, len(companies))):
            ax2.annotate(comp['name'], (comp['pe_ratio'], comp['pb_ratio']),
                         xytext=(5, 5), textcoords='offset points', fontsize=8)

        ax2.set_xlabel('PE Ratio')
        ax2.set_ylabel('PB Ratio')
        ax2.set_title('估值分布（气泡大小=市值，颜色=增速）', fontsize=12, fontweight='bold')
        ax2.grid(True, alpha=0.3)

        plt.tight_layout()

        buf = BytesIO()
        plt.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white')
        buf.seek(0)
        img = Image.open(buf)
        plt.close(fig)

        return img

    def generate_voice_transcription(self, duration_minutes: int = 30) -> Dict:
        """生成会议纪要（模拟语音转写结果）"""
        speakers = ['主持人', '研究员A', '研究员B', '行业专家', '基金经理']

        segments = []
        current_time = datetime(2024, 12, 20, 14, 0)

        # 开场白
        segments.append({
            'time': current_time.strftime('%H:%M:%S'),
            'speaker': '主持人',
            'text': '各位下午好，今天我们讨论一下新能源行业的投资机会，首先请研究员A介绍一下近期调研情况。',
            'confidence': 0.98
        })

        current_time += timedelta(minutes=2)

        # 生成对话内容
        topics = ['行业景气度', '竞争格局', '技术路线', '政策影响', '风险提示']
        for topic in topics:
            speaker = random.choice(speakers[1:])

            # 生成该话题的文本
            text = self._fallback_analysis({})
            full_text = f"关于{topic}，{text}"

            # 模拟转写特征（置信度、时间戳、语气词）
            segments.append({
                'time': current_time.strftime('%H:%M:%S'),
                'speaker': speaker,
                'text': full_text + random.choice(['', '，对吧', '，嗯', '，那个']),
                'confidence': round(random.uniform(0.85, 0.99), 2),
                'emotion': random.choice(['neutral', 'positive', 'concerned'])
            })

            # 模拟插话/打断
            if random.random() > 0.7:
                segments.append({
                    'time': (current_time + timedelta(seconds=30)).strftime('%H:%M:%S'),
                    'speaker': random.choice([s for s in speakers if s != speaker]),
                    'text': '我补充一点，' + self._fallback_analysis({}),
                    'confidence': 0.92,
                    'interruption': True
                })

            current_time += timedelta(minutes=random.randint(3, 8))

        # 总结
        segments.append({
            'time': current_time.strftime('%H:%M:%S'),
            'speaker': '主持人',
            'text': '好的，今天的讨论很有价值，后续我们会形成正式报告，散会。',
            'confidence': 0.99
        })

        return {
            'meeting_title': f'新能源行业投研讨论会 - {fake.date()}',
            'duration': f'{duration_minutes}分钟',
            'participants': list(set([s['speaker'] for s in segments])),
            'segments': segments,
            'action_items': [
                f'跟进{random.choice(self.companies)["name"]}订单情况',
                '整理政策文件',
                '下周复盘讨论'
            ]
        }


# ============================================
# 文档组装器（多模态融合）
# ============================================

class DocumentAssembler:
    """将多模态内容组装成最终文档"""

    def __init__(self, engine: MultimodalContentEngine):
        self.engine = engine
        self.output_dir = ROOT
        self.temp_files = []  # 跟踪临时文件

    def _get_temp_path(self, suffix=".png") -> str:
        """生成临时文件路径"""
        temp_path = os.path.join(TEMP_DIR, f"{uuid.uuid4().hex}{suffix}")
        self.temp_files.append(temp_path)
        return temp_path

    def _cleanup_temp_files(self):
        """清理临时文件"""
        for f in self.temp_files:
            if os.path.exists(f):
                try:
                    os.remove(f)
                except:
                    pass
        self.temp_files = []

    def create_research_report(self, filename: str) -> str:
        """生成综合研究报告（Word + 嵌入式图表）"""
        import uuid

        company = random.choice(self.engine.companies)

        # 生成AI分析文本
        analysis_text = self.engine.generate_analysis_text({
            'company': company['name'],
            'sector': company['sector'],
            'metrics': {
                'PE': company['pe_ratio'],
                'ROE': company['roe'],
                '增速': company['revenue_growth']
            }
        })

        # 创建Word文档
        doc = Document()

        # 标题
        title = doc.add_heading(f"{company['name']}深度研究报告", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # 元信息表格
        meta_table = doc.add_table(rows=1, cols=4)
        meta_table.style = 'Light Grid Accent 1'
        hdr_cells = meta_table.rows[0].cells
        headers = ['股票代码', '行业', '评级', '目标价']
        values = [company['ticker'], company['sector'],
                  random.choice(['买入', '增持', '持有']),
                  f"¥{random.randint(50, 300)}"]
        for i, (h, v) in enumerate(zip(headers, values)):
            hdr_cells[i].text = f"{h}: {v}"

        doc.add_paragraph()

        # 投资要点（带项目符号）
        doc.add_heading('投资要点', level=1)
        points = [
            f"公司是国内领先的{company['sector']}龙头，市占率持续提升",
            f"2024年预期PE {company['pe_ratio']}倍，处于历史{random.choice(['低位', '中枢', '合理区间'])}",
            f"ROE维持在{company['roe']}%水平，盈利能力稳健",
            f"{random.choice(['新产品放量', '产能扩张', '海外拓展'])}打开成长空间"
        ]
        for point in points:
            p = doc.add_paragraph(point, style='List Bullet')
            p.paragraph_format.left_indent = Inches(0.25)

        # 插入股价走势图
        doc.add_heading('股价走势分析', level=1)
        chart_img = self.engine.generate_stock_chart(company)
        chart_path = self._get_temp_path()
        chart_img.save(chart_path)
        doc.add_picture(chart_path, width=Inches(6))
        last_paragraph = doc.paragraphs[-1]
        last_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # AI分析文本
        doc.add_heading('核心逻辑', level=1)
        doc.add_paragraph(analysis_text)

        # 插入财务数据图
        doc.add_heading('财务数据', level=1)
        fin_img = self.engine.generate_financial_table_image(company)
        fin_path = self._get_temp_path()
        fin_img.save(fin_path)
        doc.add_picture(fin_path, width=Inches(6.5))

        # 估值模型表格
        doc.add_heading('盈利预测与估值', level=1)
        table = doc.add_table(rows=5, cols=4)
        table.style = 'Light Shading Accent 1'

        # 表头
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = '指标'
        hdr_cells[1].text = '2023A'
        hdr_cells[2].text = '2024E'
        hdr_cells[3].text = '2025E'

        # 数据行
        metrics_data = [
            ['营业收入(亿元)', '120.5', '145.2', '172.8'],
            ['净利润(亿元)', '15.3', '19.8', '25.4'],
            ['EPS(元)', '1.25', '1.58', '2.02'],
            ['PE(倍)', '35.2', '28.5', '22.3']
        ]
        for i, row_data in enumerate(metrics_data, 1):
            cells = table.rows[i].cells
            for j, val in enumerate(row_data):
                cells[j].text = val

        # 风险提示
        doc.add_heading('风险提示', level=1)
        risks = [
            '宏观经济下行导致需求不及预期',
            '行业竞争加剧压缩毛利率',
            '原材料价格波动影响成本',
            '政策变化带来监管风险'
        ]
        for risk in risks:
            doc.add_paragraph(risk, style='List Number')

        # 确保目录存在
        os.makedirs(os.path.dirname(filename) if os.path.dirname(filename) else self.output_dir, exist_ok=True)

        # 保存
        output_path = filename if os.path.isabs(filename) else os.path.join(self.output_dir, filename)
        doc.save(output_path)

        return output_path

    def create_presentation(self, filename: str) -> str:
        """生成路演PPT（多图多表）"""
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)

        # 标题页
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "2025年度投资策略展望"
        subtitle = title_slide.placeholders[1]
        subtitle.text = f"分析师：{fake.name()}\n日期：{datetime.now().strftime('%Y年%m月%d日')}"

        # 市场回顾页（带图表）
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])
        title1 = slide1.shapes.title
        title1.text = "市场回顾与展望"

        # 生成市场走势拼图
        fig, axes = plt.subplots(2, 2, figsize=(10, 6))

        # 指数走势
        ax1 = axes[0, 0]
        dates = self.engine.market_data['dates'][-60:]  # 最近60天
        prices = self.engine.market_data['index_prices'][-60:]
        ax1.plot(dates, prices, color=COLORS['primary'])
        ax1.set_title('沪深300走势', fontsize=10)
        ax1.grid(True, alpha=0.3)

        # 行业涨跌幅
        ax2 = axes[0, 1]
        sectors = ['科技', '医药', '消费', '金融', '新能源', '周期']
        changes = [random.uniform(-15, 25) for _ in sectors]
        colors_bar = [COLORS['positive'] if c > 0 else COLORS['negative'] for c in changes]
        ax2.barh(sectors, changes, color=colors_bar)
        ax2.set_title('行业涨跌幅(%)', fontsize=10)
        ax2.axvline(x=0, color='black', linewidth=0.5)

        # 估值分位
        ax3 = axes[1, 0]
        percentiles = [random.randint(20, 80) for _ in range(5)]
        ax3.bar(range(5), percentiles, color=COLORS['secondary'])
        ax3.set_xticks(range(5))
        ax3.set_xticklabels(['PE', 'PB', 'PS', 'PCF', 'PEG'], rotation=45)
        ax3.set_title('估值历史分位', fontsize=10)
        ax3.axhline(y=50, color='red', linestyle='--', alpha=0.5)

        # 资金流向
        ax4 = axes[1, 1]
        categories = ['北向', '公募', '私募', '保险', '外资']
        flows = [random.uniform(-50, 100) for _ in categories]
        ax4.bar(categories, flows, color=[COLORS['positive'] if f > 0 else COLORS['negative'] for f in flows])
        ax4.set_title('月度资金流向(亿元)', fontsize=10)

        plt.tight_layout()
        buf = BytesIO()
        plt.savefig(buf, format='png', dpi=150)
        buf.seek(0)
        plt.close(fig)

        # 插入PPT
        left = PptxInches(1)
        top = PptxInches(1.5)
        slide1.shapes.add_picture(buf, left, top, width=PptxInches(11))
        buf.close()

        # 个股推荐页
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        slide2.shapes.title.text = "重点推荐标的"

        # 生成对比图
        comps = random.sample(self.engine.companies, 3)
        comp_img = self.engine.generate_comparison_chart(comps)
        comp_path = self._get_temp_path()
        comp_img.save(comp_path)
        slide2.shapes.add_picture(comp_path, PptxInches(1), PptxInches(1.5), width=PptxInches(11))

        # 添加文本框说明
        left = PptxInches(1)
        top = PptxInches(6)
        width = PptxInches(11)
        height = PptxInches(1)
        textbox = slide2.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.text = "推荐逻辑：1)估值修复空间 2)业绩确定性 3)行业景气度向上"

        # 保存
        output_path = filename if os.path.isabs(filename) else os.path.join(self.output_dir, filename)
        os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else self.output_dir, exist_ok=True)
        prs.save(output_path)

        return output_path

    def create_scan_archive(self, filename: str) -> str:
        """生成扫描件档案（模拟纸质文档数字化）"""
        import uuid

        # 生成多张扫描页
        images = []
        for i in range(random.randint(3, 6)):
            img = self.engine.generate_scan_document(
                random.choice(['research_notes', 'meeting_minutes', 'draft_report'])
            )
            images.append(img)

        # 保存为PDF（模拟扫描件PDF）
        output_path = filename if os.path.isabs(filename) else os.path.join(self.output_dir, filename)
        os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else self.output_dir, exist_ok=True)

        # 使用ReportLab创建PDF
        c = SimpleDocTemplate(output_path, pagesize=A4)
        story = []

        for img in images:
            img_path = self._get_temp_path()
            img.save(img_path)
            story.append(RLImage(img_path, width=18 * cm, height=24 * cm))
            story.append(Spacer(1, 20))

        c.build(story)

        return output_path

    def create_excel_model(self, filename: str) -> str:
        """生成财务模型Excel（含公式、图表、多sheet）"""
        company = random.choice(self.engine.companies)

        wb = Workbook()

        # Sheet1: 假设条件
        ws1 = wb.active
        ws1.title = "假设条件"
        assumptions = [
            ['参数', '数值', '说明'],
            ['收入增长率(Y1)', f'{random.randint(10, 30)}%', '基于订单能见度'],
            ['毛利率', f'{random.randint(25, 45)}%', '规模效应提升'],
            ['销售费用率', f'{random.randint(5, 15)}%', '渠道扩张期'],
            ['所得税率', '15%', '高新技术企业'],
            ['WACC', f'{random.uniform(8, 12):.1f}%', '加权平均资本成本']
        ]
        for row in assumptions:
            ws1.append(row)

        # 美化
        for cell in ws1[1]:
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")

        # Sheet2: 三张表
        ws2 = wb.create_sheet("财务模型")

        # 利润表
        years = ['2022A', '2023A', '2024E', '2025E', '2026E']
        headers = ['利润表(亿元)'] + years
        ws2.append(headers)

        # 生成数据
        base_revenue = random.uniform(50, 200)
        for item in ['营业收入', '营业成本', '毛利', '销售费用', '管理费用',
                     '研发费用', '营业利润', '净利润', 'EPS']:
            row = [item]
            for i, year in enumerate(years):
                if 'A' in year:
                    # 历史数据固定
                    val = base_revenue * (0.8 if i == 0 else 1.0) * random.uniform(0.9, 1.1)
                else:
                    # 预测数据带公式（这里用值模拟）
                    growth = 1.15 if '收入' in item else 1.0
                    val = base_revenue * (growth ** (i - 1)) * random.uniform(0.95, 1.05)

                if item in ['毛利', '营业利润', '净利润']:
                    val = val * random.uniform(0.1, 0.25)  # 利润率

                row.append(round(val, 2))
            ws2.append(row)

        # 添加图表
        chart = BarChart()
        data = Reference(ws2, min_col=2, min_row=1, max_col=6, max_row=2)
        chart.add_data(data, titles_from_data=True)
        chart.title = "营收增长趋势"
        ws2.add_chart(chart, "H2")

        # Sheet3: 估值分析
        ws3 = wb.create_sheet("估值分析")
        ws3.append(['估值方法', '估值(亿元)', '权重', '加权估值'])

        methods = ['PE法', 'PB法', 'EV/EBITDA', 'DCF', '分部估值']
        total_val = 0
        weights = [0.3, 0.2, 0.2, 0.2, 0.1]

        for method, weight in zip(methods, weights):
            val = company['market_cap'] * random.uniform(0.8, 1.3)
            ws3.append([method, round(val / 1e8, 2), f'{weight:.0%}', round(val * weight / 1e8, 2)])
            total_val += val * weight

        ws3.append(['', '', '合计', round(total_val / 1e8, 2)])

        # 保存
        output_path = filename if os.path.isabs(filename) else os.path.join(self.output_dir, filename)
        os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else self.output_dir, exist_ok=True)
        wb.save(output_path)

        return output_path

    def create_meeting_archive(self, filename: str) -> str:
        """生成会议纪要（JSON格式，含转写特征）"""
        transcription = self.engine.generate_voice_transcription(duration_minutes=45)

        output_path = filename if os.path.isabs(filename) else os.path.join(self.output_dir, filename)
        os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else self.output_dir, exist_ok=True)

        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(transcription, f, ensure_ascii=False, indent=2)

        # 同时生成可读的Markdown版本
        md_content = f"""# {transcription['meeting_title']}

**时长**: {transcription['duration']}  
**参与人**: {', '.join(transcription['participants'])}

## 会议记录

| 时间 | 发言人 | 内容 | 置信度 |
|------|--------|------|--------|
"""
        for seg in transcription['segments']:
            text = seg['text'].replace('|', '\\|')
            md_content += f"| {seg['time']} | {seg['speaker']} | {text} | {seg.get('confidence', '-')} |\n"

        md_content += f"""
## 待办事项
"""
        for i, item in enumerate(transcription['action_items'], 1):
            md_content += f"{i}. {item}\n"

        md_path = output_path.replace('.json', '.md')
        with open(md_path, 'w', encoding='utf-8') as f:
            f.write(md_content)

        return output_path


# ============================================
# 批量生成主控
# ============================================

def generate_multimodal_dataset():
    """生成完整的多模态金融数据集"""
    print("🚀 初始化多模态金融数据生成引擎...")
    engine = MultimodalContentEngine()
    assembler = DocumentAssembler(engine)

    dataset_config = {
        'research_reports': {'count': 5, 'ext': 'docx', 'fn': assembler.create_research_report},
        'presentations': {'count': 3, 'ext': 'pptx', 'fn': assembler.create_presentation},
        'scan_archives': {'count': 4, 'ext': 'pdf', 'fn': assembler.create_scan_archive},
        'excel_models': {'count': 5, 'ext': 'xlsx', 'fn': assembler.create_excel_model},
        'meeting_records': {'count': 6, 'ext': 'json', 'fn': assembler.create_meeting_archive}
    }

    generated_files = []

    for category, config in dataset_config.items():
        print(f"\n📁 生成 {category} ({config['count']}个)...")
        folder = os.path.join(ROOT, category)
        os.makedirs(folder, exist_ok=True)

        for i in tqdm(range(config['count']), desc=category):
            filename = f"{category}_{i + 1:03d}.{config['ext']}"
            filepath = os.path.join(folder, filename)
            try:
                result = config['fn'](filepath)
                generated_files.append(result)
                print(f"  ✅ 已生成: {result}")
            except Exception as e:
                print(f"❌ 生成失败 {filename}: {e}")
                import traceback
                traceback.print_exc()

    # 清理临时文件
    assembler._cleanup_temp_files()

    # 生成数据集索引
    index = {
        'generated_at': datetime.now().isoformat(),
        'total_files': len(generated_files),
        'categories': {k: v['count'] for k, v in dataset_config.items()},
        'companies': [c['name'] for c in engine.companies],
        'files': generated_files
    }

    with open(os.path.join(ROOT, 'dataset_index.json'), 'w', encoding='utf-8') as f:
        json.dump(index, f, ensure_ascii=False, indent=2)

    print(f"\n✅ 多模态数据集生成完成！")
    print(f"📊 总计: {len(generated_files)} 个文件")
    print(f"📂 目录: {os.path.abspath(ROOT)}")
    print("\n包含模态类型:")
    print("  - 文本: AI生成研报、会议纪要、财务模型")
    print("  - 图表: Matplotlib生成的K线图、财务表、对比图")
    print("  - 图像: 模拟扫描件、手写笔记、截图")
    print("  - 结构化: Excel公式、JSON转写数据")
    print("  - 混合: Word/PDF中嵌入图表的复合文档")


if __name__ == "__main__":
    import uuid  # 添加uuid导入
    generate_multimodal_dataset()