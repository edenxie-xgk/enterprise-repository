import os
import random
import json
from faker import Faker
from reportlab.lib import colors
from tqdm import tqdm

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont


fake = Faker("zh_CN")

ROOT = "data"

folders = [
"01_company_profiles",
"02_industry_research",
"03_investment_memo",
"04_fund_reports",
"05_risk_management",
"06_meeting_minutes",
"07_portfolio_data",
"08_market_data",
"09_internal_training",
"10_internal_emails"
]

for f in folders:
    os.makedirs(os.path.join(ROOT,f),exist_ok=True)


# -------------------------
# 实体池
# -------------------------

industries = [
"人工智能","新能源","半导体","医疗科技",
"机器人","云计算","消费科技"
]

funds = [
"HC Alpha Growth Fund",
"HC AI Innovation Fund",
"HC Global Macro Fund",
"HC Energy Transition Fund",
"HC Quant Arbitrage Fund"
]

departments = [
"投资策略部",
"量化研究部",
"风险管理部",
"资产配置委员会",
"投资委员会"
]


# -------------------------
# 生成公司数据库
# -------------------------

companies = []

for i in range(40):

    company = {
        "name":fake.company(),
        "industry":random.choice(industries),
        "founded":random.randint(2005,2020),
        "employees":random.randint(200,8000),
        "revenue":random.randint(10,500)
    }

    companies.append(company)



# -------------------------
# 文本模板池
# -------------------------

paragraph_templates = [

"""
{company} 是中国 {industry} 行业的重要企业之一。
根据华辰资本 {dept} 的研究，该公司成立于 {founded} 年，
目前员工规模约 {employees} 人。公司近年来持续加大研发投入，
在核心技术方面形成了较强的竞争壁垒。
""",

"""
从产业链角度看，{industry} 行业主要由上游设备厂商、
中游平台企业以及下游应用厂商组成。
{company} 主要位于产业链的 {position} 环节。
""",

"""
华辰资本旗下的 {fund} 在最近的资产配置会议中
讨论了 {industry} 行业的投资机会，并对 {company}
进行了重点评估。
""",

"""
风险管理部认为，{industry} 行业仍然面临若干
不确定性因素，包括技术迭代风险、政策变化、
以及市场竞争加剧等问题。
""",

"""
根据内部财务模型测算，{company} 在未来三年
收入复合增长率预计可达到 {growth}% 左右。
该预测基于行业需求增长以及公司产能扩张计划。
"""
]


def generate_paragraph():

    template=random.choice(paragraph_templates)

    company=random.choice(companies)

    return template.format(
        company=company["name"],
        industry=company["industry"],
        dept=random.choice(departments),
        fund=random.choice(funds),
        founded=company["founded"],
        employees=company["employees"],
        growth=random.randint(10,35),
        position=random.choice(["上游","中游","下游"])
    )


# -------------------------
# Markdown
# -------------------------

def generate_md(path):

    with open(path,"w",encoding="utf8") as f:

        f.write("# 行业深度研究报告\n\n")

        for i in range(5):

            f.write(f"## 行业分析章节 {i+1}\n\n")

            for _ in range(10):

                f.write(generate_paragraph()+"\n\n")

            f.write("| 公司 | 行业 | 市值 | 增长率 |\n")
            f.write("|----|----|----|----|\n")

            for _ in range(20):

                c=random.choice(companies)

                f.write(
                    f"| {c['name']} | {c['industry']} | {random.randint(100,800)}亿 | {random.randint(5,30)}% |\n"
                )

            f.write("\n\n")


# -------------------------
# DOCX
# -------------------------

def generate_docx(path):

    doc=Document()

    doc.add_heading("投资研究报告",0)

    for i in range(5):

        doc.add_heading(f"行业分析 {i+1}",level=1)

        for _ in range(8):
            doc.add_paragraph(generate_paragraph())

        table=doc.add_table(rows=1,cols=4)

        hdr=table.rows[0].cells
        hdr[0].text="公司"
        hdr[1].text="行业"
        hdr[2].text="市值"
        hdr[3].text="增长率"

        for _ in range(15):

            c=random.choice(companies)

            row=table.add_row().cells
            row[0].text=c["name"]
            row[1].text=c["industry"]
            row[2].text=str(random.randint(100,800))+"亿"
            row[3].text=str(random.randint(5,30))+"%"

    doc.save(path)



# -------------------------
# PDF
# -------------------------


style = ParagraphStyle(
    "CN",
    fontName="NotoSans",
    fontSize=11,
    leading=16
)
FONT_PATH = "fonts/NotoSansCJKsc-VF.ttf"

pdfmetrics.registerFont(
    TTFont("NotoSans", FONT_PATH)
)

def generate_pdf(path):

    story=[]

    story.append(Paragraph("基金投资研究报告",style))
    story.append(Spacer(1,20))

    for _ in range(6):

        story.append(Paragraph(generate_paragraph(),style))
        story.append(Spacer(1,10))

    table_data=[["公司","行业","市值","增长率"]]

    for _ in range(20):

        c=random.choice(companies)

        table_data.append([
            c["name"],
            c["industry"],
            str(random.randint(100,800))+"亿",
            str(random.randint(5,30))+"%"
        ])
    story.append(Paragraph("主要投资结论", style))

    story.append(Paragraph(
        "• 人工智能行业未来五年预计保持15%增长<br/>"
        "• 星云科技集团具备较强技术壁垒<br/>"
        "• HC AI Innovation Fund 已进行战略投资",
        style))

    table=Table(table_data,colWidths=[6*cm,4*cm,3*cm,3*cm])
    table.setStyle(TableStyle([
        ('FONTNAME', (0, 0), (-1, -1), 'NotoSans'),
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('ALIGN', (2, 1), (-1, -1), 'RIGHT'),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
        ('BACKGROUND', (0, 0), (-1, 0), colors.lightgrey)
    ]))
    story.append(table)
    story.append(Spacer(1, 20))
    pdf=SimpleDocTemplate(path)

    pdf.build(story)



# -------------------------
# Excel
# -------------------------

def generate_excel(path):

    wb=Workbook()

    ws=wb.active
    ws.title="基金持仓"

    ws.append(["公司","行业","持仓金额","收益率"])

    for _ in range(500):

        c=random.choice(companies)

        ws.append([
            c["name"],
            c["industry"],
            random.randint(100000000,5000000000),
            round(random.uniform(-5,25),2)
        ])

    wb.save(path)



# -------------------------
# JSON
# -------------------------

def generate_json(path):

    data=[]

    for _ in range(3000):

        c=random.choice(companies)

        item={
            "date":fake.date(),
            "company":c["name"],
            "industry":c["industry"],
            "price":round(random.uniform(20,200),2),
            "volume":random.randint(100000,5000000)
        }

        data.append(item)

    with open(path,"w",encoding="utf8") as f:
        json.dump(data,f,ensure_ascii=False,indent=2)



# -------------------------
# PPT
# -------------------------

def generate_ppt(path):

    prs=Presentation()

    for i in range(15):

        slide_layout=prs.slide_layouts[1]

        slide=prs.slides.add_slide(slide_layout)

        slide.shapes.title.text=f"投资策略分析 {i+1}"

        content=slide.placeholders[1]

        text="\n".join(generate_paragraph() for _ in range(3))

        content.text=text

    prs.save(path)



# -------------------------
# TXT
# -------------------------

def generate_txt(path):

    with open(path,"w",encoding="utf8") as f:

        for _ in range(200):

            f.write(generate_paragraph())
            f.write("\n\n")



# -------------------------
# 生成全部文件
# -------------------------

def generate_all():

    for i in tqdm(range(60)):
        generate_md(f"{ROOT}/02_industry_research/industry_report_{i+1}.md")

    for i in tqdm(range(40)):
        generate_docx(f"{ROOT}/03_investment_memo/investment_memo_{i+1}.docx")

    for i in tqdm(range(35)):
        generate_pdf(f"{ROOT}/04_fund_reports/fund_report_{i+1}.pdf")

    for i in tqdm(range(20)):
        generate_excel(f"{ROOT}/07_portfolio_data/portfolio_data_{i+1}.xlsx")

    for i in tqdm(range(20)):
        generate_json(f"{ROOT}/08_market_data/market_data_{i+1}.json")

    for i in tqdm(range(15)):
        generate_ppt(f"{ROOT}/09_internal_training/training_{i+1}.pptx")

    for i in tqdm(range(15)):
        generate_txt(f"{ROOT}/10_internal_emails/internal_mail_{i+1}.txt")



if __name__=="__main__":

    generate_all()

    print("金融RAG知识库生成完成")